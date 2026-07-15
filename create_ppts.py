"""
Generate 3 PowerPoint presentations for Layers 1, 2, and 3
of the Human Microbiome Research Knowledge Graph pipeline.
Run: python3 create_ppts.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ────────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x0D, 0x1B, 0x2A)
CARD_BG    = RGBColor(0x16, 0x2A, 0x3E)
ACCENT1    = RGBColor(0x00, 0xB4, 0xD8)   # cyan        – Layer 1
ACCENT2    = RGBColor(0x48, 0xCA, 0x9A)   # teal-green  – Layer 2
ACCENT3    = RGBColor(0xF7, 0xC5, 0x9F)   # warm amber  – Layer 3
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xD6, 0xE0)

W = Inches(13.33)
H = Inches(7.5)

# ── Primitive helpers ──────────────────────────────────────────────────────
def new_prs():
    p = Presentation()
    p.slide_width  = W
    p.slide_height = H
    return p

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color=DARK_BG):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def box(slide, x, y, w, h, fill=None, border=None, bpt=0):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.line.width = Pt(bpt)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border and bpt > 0:
        s.line.color.rgb = border
    else:
        s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h, size=16, bold=False,
        color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text            = text
    r.font.size       = Pt(size)
    r.font.bold       = bold
    r.font.color.rgb  = color
    return tb

def top_bar(slide, color):
    box(slide, 0, 0, W, Inches(0.07), fill=color)

def bot_bar(slide, color):
    box(slide, 0, H - Inches(0.07), W, Inches(0.07), fill=color)

def slide_num(slide, n, total, color=LIGHT_GRAY):
    txt(slide, f"{n} / {total}",
        W - Inches(1.2), H - Inches(0.38),
        Inches(1.1), Inches(0.32),
        size=11, color=color, align=PP_ALIGN.RIGHT)

def section_header(slide, title, subtitle, accent):
    bg(slide)
    top_bar(slide, accent); bot_bar(slide, accent)
    txt(slide, title,
        Inches(0.6), Inches(2.6), Inches(12), Inches(1.0),
        size=48, bold=True, color=WHITE)
    txt(slide, subtitle,
        Inches(0.6), Inches(3.8), Inches(12), Inches(0.6),
        size=20, color=accent)

def cards(slide, items, start_y, cols, card_w, card_h, gap_x=Inches(0.18),
          gap_y=Inches(0.22)):
    """
    items: list of (title, body, accent_color)
    """
    for i, (title, body, color) in enumerate(items):
        col = i % cols
        row = i // cols
        cx  = Inches(0.5) + col * (card_w + gap_x)
        cy  = start_y      + row * (card_h + gap_y)
        box(slide, cx, cy, card_w, card_h,
            fill=CARD_BG, border=color, bpt=2)
        txt(slide, title,
            cx + Inches(0.14), cy + Inches(0.1),
            card_w - Inches(0.28), Inches(0.42),
            size=13, bold=True, color=color)
        txt(slide, body,
            cx + Inches(0.14), cy + Inches(0.55),
            card_w - Inches(0.28), card_h - Inches(0.65),
            size=12, color=LIGHT_GRAY)

def two_col(slide, left_title, left_lines, right_title, right_lines,
            left_color, right_color, top=Inches(1.1)):
    lw = Inches(5.9); rw = Inches(6.0)
    lx = Inches(0.5); rx = Inches(6.8)
    bh = Inches(5.6)
    box(slide, lx, top, lw, bh, fill=CARD_BG, border=left_color, bpt=1)
    txt(slide, left_title, lx+Inches(0.15), top+Inches(0.1),
        lw-Inches(0.3), Inches(0.42), size=16, bold=True, color=left_color)
    for i, line in enumerate(left_lines):
        txt(slide, line, lx+Inches(0.15), top+Inches(0.65)+Inches(0.42*i),
            lw-Inches(0.3), Inches(0.4), size=13, color=LIGHT_GRAY)
    box(slide, rx, top, rw, bh, fill=CARD_BG, border=right_color, bpt=1)
    txt(slide, right_title, rx+Inches(0.15), top+Inches(0.1),
        rw-Inches(0.3), Inches(0.42), size=16, bold=True, color=right_color)
    for i, line in enumerate(right_lines):
        txt(slide, line, rx+Inches(0.15), top+Inches(0.65)+Inches(0.42*i),
            rw-Inches(0.3), Inches(0.4), size=13, color=LIGHT_GRAY)

# ══════════════════════════════════════════════════════════════════════════
#  LAYER 1 — Data Collection  (10 slides)
# ══════════════════════════════════════════════════════════════════════════
def make_layer1():
    prs = new_prs(); A = ACCENT1; T = 10

    # 1 — Title
    s = blank(prs); bg(s); top_bar(s, A); bot_bar(s, A)
    txt(s, "Layer 1", Inches(0.6), Inches(1.1), Inches(12), Inches(0.7),
        size=22, bold=True, color=A)
    txt(s, "Data Collection Pipeline",
        Inches(0.6), Inches(1.9), Inches(12), Inches(1.3),
        size=50, bold=True, color=WHITE)
    txt(s, "Human Microbiome Research Knowledge Graph",
        Inches(0.6), Inches(3.4), Inches(12), Inches(0.6),
        size=22, color=LIGHT_GRAY)
    txt(s, "6 APIs  ·  4-Stage Filter  ·  Smart Dedup  ·  Cursor Persistence",
        Inches(0.6), Inches(4.1), Inches(12), Inches(0.5),
        size=18, color=A)
    slide_num(s, 1, T)

    # 2 — Pipeline position
    s = blank(prs); bg(s); top_bar(s, A)
    txt(s, "Where Layer 1 Fits",
        Inches(0.6), Inches(0.18), Inches(10), Inches(0.55),
        size=26, bold=True, color=A)
    stages = [
        ("LAYER 1  ◀ YOU ARE HERE", "Fetch · Dedup · Filter", A, DARK_BG),
        ("LAYER 2",  "NLP Enrichment",          ACCENT2, DARK_BG),
        ("LAYER 3",  "Knowledge Graph (Neo4j)", ACCENT3, DARK_BG),
        ("QUERIES",  "REST API · 5 Research Qs",LIGHT_GRAY, DARK_BG),
    ]
    for i, (ti, su, fc, _) in enumerate(stages):
        bx = Inches(0.5 + i * 3.18)
        box(s, bx, Inches(1.1), Inches(3.0), Inches(1.0),
            fill=CARD_BG, border=fc, bpt=2)
        txt(s, ti, bx+Inches(0.12), Inches(1.15), Inches(2.76),
            Inches(0.4), size=12, bold=True, color=fc)
        txt(s, su, bx+Inches(0.12), Inches(1.58), Inches(2.76),
            Inches(0.4), size=12, color=LIGHT_GRAY)
        if i < 3:
            txt(s, "▶", bx + Inches(3.02), Inches(1.48),
                Inches(0.2), Inches(0.4), size=18, color=LIGHT_GRAY)
    bullets = [
        "• 6 collectors run in parallel — PubMed · EuropePMC · SemanticScholar · OpenAlex · Crossref · CORE",
        "• Results merged and deduplicated: DOI → PMID → title[:80]",
        "• Best field value from each source kept on merge (citation count from S2, MeSH from PubMed…)",
        "• 4-stage relevance filter applied to every paper before saving",
        "• Cursor state saved after each run — next run resumes from exact offset",
        "• Output:  data/processed/collected_YYYYMMDD_HHMMSS.json  →  List[PaperRecord]",
    ]
    for i, b in enumerate(bullets):
        txt(s, b, Inches(0.6), Inches(2.35)+Inches(0.55*i),
            Inches(12.2), Inches(0.5), size=15, color=LIGHT_GRAY)
    slide_num(s, 2, T)

    # 3 — 6 Data Sources
    s = blank(prs); bg(s); top_bar(s, A)
    txt(s, "6 Data Sources — Why Each One Matters",
        Inches(0.6), Inches(0.18), Inches(12), Inches(0.55),
        size=26, bold=True, color=A)
    sources = [
        ("PubMed",           "MeSH terms · Human filter\nGold-standard biomedical index\n3 req/s → 10 req/s with API key", A),
        ("Europe PMC",       "Full-text XML · PMCIDs\nmedRxiv preprint coverage\nClean JSON API", ACCENT2),
        ("Semantic Scholar", "Citation counts · Bulk endpoint\nToken pagination · AI field tags\n1 req/s with API key", ACCENT3),
        ("OpenAlex",         "250M+ works · Concept IDs\nFunding metadata · Cursor pagination\nPolite pool 100 req/s", A),
        ("Crossref",         "Canonical DOIs · Funder names\nJATS XML abstract stripping\nOffset pagination", ACCENT2),
        ("CORE",             "Full-text PDFs · OA aggregator\nActual paper content for NLP\n1 000 tokens/day with key", ACCENT3),
    ]
    cards(s, sources, Inches(0.95), 3, Inches(4.08), Inches(1.55))
    txt(s, "All 6 run in parallel via ThreadPoolExecutor — independent API hosts, no shared bottleneck.",
        Inches(0.6), Inches(7.05), Inches(12.2), Inches(0.35),
        size=13, color=A, align=PP_ALIGN.CENTER)
    slide_num(s, 3, T)

    # 4 — PaperRecord data model
    s = blank(prs); bg(s); top_bar(s, A)
    txt(s, "PaperRecord — Shared Contract Between All Layers",
        Inches(0.6), Inches(0.18), Inches(12), Inches(0.55),
        size=26, bold=True, color=A)
    txt(s, "Every collector converts its raw API response into one standardised PaperRecord. "
           "Layer 2 never needs to know which source a paper came from.",
        Inches(0.6), Inches(0.82), Inches(12.2), Inches(0.5), size=15, color=LIGHT_GRAY)
    groups = [
        ("Identity",          "doi · pmid · pmcid · arxiv_id · source",              A),
        ("Core Content",      "title · abstract · authors · keywords",               ACCENT2),
        ("Publication Info",  "journal · issn · pub_date · volume · issue · pages",  ACCENT3),
        ("Classification",    "article_types  e.g. [Journal Article, RCT]",         A),
        ("Access & Full Text","is_open_access · full_text_url · pdf_url",           ACCENT2),
        ("Citations",         "citation_count · reference_count",                    ACCENT3),
        ("PubMed-specific",   "mesh_terms  (key for Stage 1 filter)",               A),
        ("Pipeline Meta",     "content_hash · fetched_at · is_preprint",            ACCENT2),
    ]
    cards(s, groups, Inches(1.45), 2, Inches(5.9), Inches(0.75), gap_y=Inches(0.1))
    txt(s, "Dedup key:  DOI  →  PMID  →  title[:80]  (fuzzy fallback)",
        Inches(0.6), Inches(6.95), Inches(12), Inches(0.38),
        size=13, color=ACCENT2)
    slide_num(s, 4, T)

    # 5 — 4-Stage Filter
    s = blank(prs); bg(s); top_bar(s, A)
    txt(s, "4-Stage Relevance Filter — Cheapest Stage First",
        Inches(0.6), Inches(0.18), Inches(12), Inches(0.55),
        size=26, bold=True, color=A)
    stages_f = [
        ("Stage 1", "MeSH Metadata",
         "Zero cost — dictionary lookups\nHas human + microbiome MeSH → KEEP\nHas animal-only MeSH → REJECT\nNo MeSH → borderline → Stage 2", A),
        ("Stage 2", "Weighted Keywords",
         "Score over title+abstract+mesh\n≥0.70 → KEEP → Metagenomics Gate\n0.40–0.69 → borderline → Stage 3\n<0.40 → REJECT", ACCENT2),
        ("Stage 3", "ML Classifier",
         "sentence-transformers + LogisticRegression\n384-dim embedding → probability\nSelf-supervised bootstrap training\nBorderline → Stage 3.5 → Stage 4", ACCENT3),
        ("Stage 4", "LLM Verifier",
         "Local Ollama llama3 — borderlines only\nFull title+abstract reasoning\nSemantic cache (0.97 similarity reuse)\nBatched: 16 papers per LLM call", LIGHT_GRAY),
    ]
    cw = Inches(3.0); top = Inches(1.0)
    for i, (num, name, desc, color) in enumerate(stages_f):
        cx = Inches(0.5) + i*(cw+Inches(0.11))
        box(s, cx, top, cw, Inches(5.95), fill=CARD_BG, border=color, bpt=2)
        txt(s, num,  cx+Inches(0.14), top+Inches(0.1), cw-Inches(0.28),
            Inches(0.34), size=12, bold=True, color=color)
        txt(s, name, cx+Inches(0.14), top+Inches(0.48), cw-Inches(0.28),
            Inches(0.46), size=16, bold=True, color=WHITE)
        txt(s, desc, cx+Inches(0.14), top+Inches(1.02), cw-Inches(0.28),
            Inches(4.7), size=12, color=LIGHT_GRAY)
        if i < 3:
            txt(s, "▶", cx+cw+Inches(0.0), top+Inches(2.8),
                Inches(0.14), Inches(0.45), size=18, color=LIGHT_GRAY)
    txt(s, "Metagenomics Gate (after Stage 2 KEEP): paper must mention ≥1 of 200+ sequencing/data terms.",
        Inches(0.6), Inches(7.1), Inches(12.2), Inches(0.32),
        size=12, color=ACCENT2)
    slide_num(s, 5, T)

    # 6 — Orchestrator & Deduplication
    s = blank(prs); bg(s); top_bar(s, A)
    txt(s, "Orchestrator — Parallel Collection & Smart Deduplication",
        Inches(0.6), Inches(0.18), Inches(12), Inches(0.55),
        size=26, bold=True, color=A)
    two_col(s,
        "Parallel Collection",
        ["• ThreadPoolExecutor — 6 workers",
         "• Each hits a different API host",
         "• GIL released during HTTP I/O",
         "• Threads block in parallel (true concurrency)",
         "",
         "Rate limits per collector:",
         "PubMed  0.1s  |  OpenAlex  0.1s",
         "Crossref 0.02s  |  CORE  0.6s",
         "EuropePMC 0.5s  |  S2  1.0s"],
        "Deduplication & Merging",
        ["Priority:  DOI → PMID → title[:80]",
         "",
         "Same paper in N sources? → MERGE:",
         "  PubMed     wins  mesh_terms, pub_date",
         "  S2         wins  citation_count",
         "  EuropePMC  wins  pmcid, full_text_url",
         "  Crossref   wins  doi, journal_abbrev",
         "",
         "source = 'merged:pubmed+europepmc+…'",
         "",
         "Cursors saved → next run resumes from",
         "exact offset / token / cursor"],
        A, ACCENT2)
    slide_num(s, 6, T)

    # 7 — BaseCollector Infrastructure
    s = blank(prs); bg(s); top_bar(s, A)
    txt(s, "BaseCollector — Shared Infrastructure",
        Inches(0.6), Inches(0.18), Inches(12), Inches(0.55),
        size=26, bold=True, color=A)
    infra = [
        ("Rate Limiting",     "Configurable min-gap per source.\nPrevents API bans automatically.", A),
        ("Auto Retry",        "tenacity exponential backoff:\n2 s → 4 s → 8 s → give up.", ACCENT2),
        ("Content Hash",      "MD5(title+abstract) detects corrections.\nSkips unchanged re-fetches.", ACCENT3),
        ("Raw Cache",         "Every API response saved to\ndata/raw/<source>/ for replay.", A),
        ("User-Agent ID",     "Academic-friendly identification.\nAnonymous scrapers get blocked.", ACCENT2),
        ("PubMed WebHistory", "2-step esearch+efetch.\nPages all 70K+ monthly results.", ACCENT3),
    ]
    cards(s, infra, Inches(0.95), 3, Inches(4.08), Inches(1.55))
    txt(s, "Each collector implements only 3 methods:  build_query()  ·  fetch_page()  ·  parse_record()",
        Inches(0.6), Inches(7.05), Inches(12.2), Inches(0.35),
        size=13, color=ACCENT2, align=PP_ALIGN.CENTER)
    slide_num(s, 7, T)

    # 8 — Output & Audit
    s = blank(prs); bg(s); top_bar(s, A)
    txt(s, "Output, Audit & Observability",
        Inches(0.6), Inches(0.18), Inches(12), Inches(0.55),
        size=26, bold=True, color=A)
    two_col(s,
        "Output File Structure",
        ["data/processed/",
         "  collected_YYYYMMDD_HHMMSS.json   ← Layer 1 output",
         "  collector_cursors.json           ← cross-run state",
         "  llm_cache.json                   ← Stage 4 cache",
         "data/raw/<source>/                 ← API response cache",
         "data/embeddings/",
         "  positive_embeddings.npy          ← Stage 3.5",
         "  negative_embeddings.npy",
         "data/audit/",
         "  kept.json   rejected.json",
         "  review.json  llm_verified.json"],
        "Audit Trail per Decision",
        ["Every keep / reject / review logged:",
         "",
         '  "title": "Gut microbiome in IBD…",',
         '  "source": "pubmed",',
         '  "decision": "keep",',
         '  "stage": "stage1_metadata",',
         '  "score": 0.90,',
         '  "reason": "human+microbiome_mesh",',
         '  "doi": "10.1038/…"',
         "",
         "Fully auditable: trace any paper's fate.",
         "Stage 2 Calibrator auto-tunes thresholds."],
        A, ACCENT2)
    slide_num(s, 8, T)

    # 9 — Cross-run incrementality
    s = blank(prs); bg(s); top_bar(s, A)
    txt(s, "Cross-Run Incrementality & Configuration",
        Inches(0.6), Inches(0.18), Inches(12), Inches(0.55),
        size=26, bold=True, color=A)
    two_col(s,
        "Cursor Persistence",
        ["collector_cursors.json tracks position:",
         "",
         '  "pubmed": 15000,',
         '  "europepmc": 3000,',
         '  "semantic_scholar_token": "eyJ…",',
         '  "openalex_cursor": "IljD…",',
         '  "crossref": 2000,',
         '  "core": 500',
         "",
         "Delete file → re-collect everything.",
         "Each source resumes exact offset,",
         "token, or cursor independently."],
        "Key .env Settings",
        ["NCBI_EMAIL            required",
         "NCBI_API_KEY          10 req/s",
         "SEMANTIC_SCHOLAR_API_KEY",
         "CORE_API_KEY          1 000 tokens/day",
         "",
         "MAX_RESULTS_PER_SOURCE   (default 500)",
         "METAGENOMICS_GATE_ENABLED  (default true)",
         "OLLAMA_VERIFIER_MODEL    (default llama3)",
         "EMBEDDING_MODEL_NAME     (SPECTER2)",
         "SEMANTIC_CACHE_THRESHOLD (0.97)",
         "BATCH_LLM_SIZE           (16)"],
        A, ACCENT2)
    slide_num(s, 9, T)

    # 10 — Summary
    s = blank(prs); bg(s); top_bar(s, A); bot_bar(s, A)
    txt(s, "Layer 1 — Summary",
        Inches(0.6), Inches(0.22), Inches(12), Inches(0.55),
        size=28, bold=True, color=A)
    summary = [
        ("6 Sources",     "PubMed · EuropePMC · S2\nOpenAlex · Crossref · CORE",  A),
        ("4-Stage Filter","MeSH → Keywords → ML → LLM\nCheapest stage first",      ACCENT2),
        ("Smart Dedup",   "DOI → PMID → title merge\nBest field from each source", ACCENT3),
        ("Full Audit",    "Every decision logged with\nscore, stage, reason, DOI", A),
        ("Incremental",   "Cursors persist across runs\nReset by deleting file",   ACCENT2),
        ("Output",        "List[PaperRecord] →\ncollected_YYYYMMDD.json",          ACCENT3),
    ]
    cards(s, summary, Inches(1.1), 3, Inches(4.08), Inches(1.6))
    slide_num(s, 10, T)

    prs.save("Layer1_DataCollection.pptx")
    print("✓  Layer1_DataCollection.pptx")

# ══════════════════════════════════════════════════════════════════════════
#  LAYER 2 — NLP Enrichment  (10 slides)
# ══════════════════════════════════════════════════════════════════════════
def make_layer2():
    prs = new_prs(); A = ACCENT2; T = 10

    # 1 — Title
    s = blank(prs); bg(s); top_bar(s, A); bot_bar(s, A)
    txt(s, "Layer 2", Inches(0.6), Inches(1.1), Inches(12), Inches(0.7),
        size=22, bold=True, color=A)
    txt(s, "NLP Enrichment Pipeline",
        Inches(0.6), Inches(1.9), Inches(12), Inches(1.3),
        size=48, bold=True, color=WHITE)
    txt(s, "Human Microbiome Research Knowledge Graph",
        Inches(0.6), Inches(3.4), Inches(12), Inches(0.6),
        size=22, color=LIGHT_GRAY)
    txt(s, "Entities · Classification · Full Text · Quality · Ontology Grounding",
        Inches(0.6), Inches(4.1), Inches(12), Inches(0.5),
        size=18, color=A)
    slide_num(s, 1, T)

    # 2 — 4-Phase Pipeline Overview
    s = blank(prs); bg(s); top_bar(s, A)
    txt(s, "4-Phase Pipeline Overview",
        Inches(0.6), Inches(0.18), Inches(12), Inches(0.55),
        size=26, bold=True, color=A)
    phases = [
        ("Phase 0",   "Batch PMCID\nResolution",
         "200 DOIs per NCBI request\nEuropePMC fallback\n10× quality improvement", A),
        ("Phase 0.5", "PMC Full-Text\nPre-Enrichment",
         "Pre-fetch structured XML\nbefore Phase 1 threads start\nAvoids race on PMC API", ACCENT2),
        ("Phase 1",   "Full-Text\nAcquisition",
         "64 parallel I/O threads\nTiered: PMC → PDF → Web\n→ Abstract → exhausted", ACCENT3),
        ("Phase 2",   "NLP Processing",
         "4–8 CPU processes\n8 modules per paper:\nNER · classify · score", LIGHT_GRAY),
    ]
    pw = Inches(3.0); top = Inches(1.0)
    for i, (num, name, desc, color) in enumerate(phases):
        px = Inches(0.5) + i*(pw+Inches(0.11))
        box(s, px, top, pw, Inches(5.95), fill=CARD_BG, border=color, bpt=2)
        txt(s, num,  px+Inches(0.14), top+Inches(0.1), pw-Inches(0.28),
            Inches(0.34), size=12, bold=True, color=color)
        txt(s, name, px+Inches(0.14), top+Inches(0.48), pw-Inches(0.28),
            Inches(0.7), size=16, bold=True, color=WHITE)
        txt(s, desc, px+Inches(0.14), top+Inches(1.25), pw-Inches(0.28),
            Inches(4.5), size=13, color=LIGHT_GRAY)
        if i < 3:
            txt(s, "▶", px+pw+Inches(0.0), top+Inches(2.8),
                Inches(0.14), Inches(0.45), size=18, color=LIGHT_GRAY)
    slide_num(s, 2, T)

    # 3 — Full-Text Acquisition Tiers
    s = blank(prs); bg(s); top_bar(s, A)
    txt(s, "Phase 1 — Full-Text Acquisition (64 I/O Threads)",
        Inches(0.6), Inches(0.18), Inches(12), Inches(0.55),
        size=26, bold=True, color=A)
    txt(s, "FullTextOrchestrator tries strategies in priority order — always highest-quality source first.",
        Inches(0.6), Inches(0.82), Inches(12.2), Inches(0.42), size=15, color=LIGHT_GRAY)
    tiers = [
        ("Tier 1 — Structured XML", "Has PMCID?",
         "EuropePMC XML  →  NCBI PMC XML\nLabeled sections: Methods / Results / Data Availability\nFastest. Best NLP quality.", A),
        ("Tier 2 — OA Content",     "Has PDF / URL / DOI?",
         "PDF Parser (pymupdf4llm + OCR fallback)\nWeb Scraper (trafilatura HTML extraction)\nUnpaywall → OpenAIRE (EU repositories)", ACCENT2),
        ("Tier 3 — Structured Abstract", "Has PMID?",
         "NCBI efetch: complete structured abstract\nBackground / Methods / Results / Conclusions\nHigher quality than collector's raw abstract.", ACCENT3),
        ("Tier 4 — Fallback",       "Nothing worked",
         "Mark exhausted in fetch_cache.json\nRetry after 90 days (OA embargoes lift)\nUse collector's abstract for NLP.", LIGHT_GRAY),
    ]
    tw = Inches(5.95); th = Inches(1.48)
    for i, (tier, cond, desc, color) in enumerate(tiers):
        col = i % 2; row = i // 2
        tx = Inches(0.5) + col*(tw+Inches(0.35))
        ty = Inches(1.55) + row*(th+Inches(0.22))
        box(s, tx, ty, tw, th, fill=CARD_BG, border=color, bpt=2)
        txt(s, tier, tx+Inches(0.14), ty+Inches(0.07), tw-Inches(0.28),
            Inches(0.36), size=13, bold=True, color=color)
        txt(s, cond, tx+Inches(0.14), ty+Inches(0.45), tw-Inches(0.28),
            Inches(0.28), size=12, bold=True, color=WHITE)
        txt(s, desc, tx+Inches(0.14), ty+Inches(0.75), tw-Inches(0.28),
            Inches(0.65), size=12, color=LIGHT_GRAY)
    txt(s, "Cache:  fetch_cache.json  keyed by content_hash.  Success → instant replay.  "
           "Exhausted → skip 90 days.",
        Inches(0.6), Inches(7.1), Inches(12.2), Inches(0.32),
        size=12, color=ACCENT2)
    slide_num(s, 3, T)

    # 4 — 8 NLP Modules
    s = blank(prs); bg(s); top_bar(s, A)
    txt(s, "Phase 2 — 8 NLP Modules (4–8 CPU Processes)",
        Inches(0.6), Inches(0.18), Inches(12), Inches(0.55),
        size=26, bold=True, color=A)
    modules = [
        ("1", "Article Classifier",   "meta_analysis / rct / review\narticle_type_confidence 0–1", A),
        ("2", "Journal Classifier",   "IF · Q1–Q4 quartile · OA flag\nPredatory publisher check",  ACCENT2),
        ("3", "Section Parser",       "26 section types: methods,\nresults, data_availability…",   ACCENT3),
        ("4", "NER Extractor",        "3-tier: Regex → BioBERT → LLM\n18 entity categories",       LIGHT_GRAY),
        ("5", "Data Availability",    "SRA / GEO accession IDs\nopen / restricted / not_stated",   A),
        ("6", "Study Design",         "RCT · cohort · case-control\nCalibrated confidence per type",ACCENT2),
        ("7", "Evidence Extractor",   "Sample size · seq methods\nDataset accession IDs",          ACCENT3),
        ("8", "Quality Scorer",       "Composite 0.0–1.0:\nJournal + Design + Sample + Data",      LIGHT_GRAY),
    ]
    mw = Inches(2.95); mh = Inches(1.45)
    for i, (num, name, desc, color) in enumerate(modules):
        col = i % 4; row = i // 4
        mx = Inches(0.5) + col*(mw+Inches(0.1))
        my = Inches(1.0) + row*(mh+Inches(0.22))
        box(s, mx, my, mw, mh, fill=CARD_BG, border=color, bpt=2)
        txt(s, f"Module {num}", mx+Inches(0.12), my+Inches(0.05),
            mw-Inches(0.24), Inches(0.3), size=11, color=color)
        txt(s, name, mx+Inches(0.12), my+Inches(0.38),
            mw-Inches(0.24), Inches(0.4), size=14, bold=True, color=WHITE)
        txt(s, desc, mx+Inches(0.12), my+Inches(0.82),
            mw-Inches(0.24), Inches(0.55), size=12, color=LIGHT_GRAY)
    slide_num(s, 4, T)

    # 5 — 3-Tier NER
    s = blank(prs); bg(s); top_bar(s, A)
    txt(s, "Module 4: 3-Tier NER — Named Entity Recognition",
        Inches(0.6), Inches(0.18), Inches(12), Inches(0.55),
        size=26, bold=True, color=A)
    tiers_ner = [
        ("Tier 1", "Rule-Based Regex", "ALWAYS ON",
         "18 curated category dictionaries\n~300 compiled patterns (longest-first)\n~5 ms/paper  ·  Precision 95%+\nAbbreviation confirmation via context window", A, "~5 ms/paper"),
        ("Tier 2", "BioBERT NER Model", "use_model=True",
         "d4data/biomedical-ner-all  (440 MB)\nCRAFT + BioNLP pre-trained corpora\n400-word chunks, 50-word overlap\nGPU: ~50 ms  ·  CPU: ~500 ms/paper", ACCENT2, "~500 ms/paper CPU"),
        ("Tier 3", "LLM Extraction", "use_llm=True",
         "Ollama llama3 (cross-process FileLock)\nGap-filling: entities Tier 1+2 missed\nTop-priority sections up to 3 000 chars\nResults cached in triple_cache.json", ACCENT3, "~30–120 s/paper"),
    ]
    tw2 = Inches(3.95); th2 = Inches(5.5); top2 = Inches(1.0)
    for i, (num, name, flag, desc, color, speed) in enumerate(tiers_ner):
        tx2 = Inches(0.5) + i*(tw2+Inches(0.25))
        box(s, tx2, top2, tw2, th2, fill=CARD_BG, border=color, bpt=2)
        txt(s, num,  tx2+Inches(0.14), top2+Inches(0.1), tw2-Inches(0.28),
            Inches(0.3), size=12, bold=True, color=color)
        txt(s, name, tx2+Inches(0.14), top2+Inches(0.42), tw2-Inches(0.28),
            Inches(0.46), size=16, bold=True, color=WHITE)
        txt(s, f"({flag})", tx2+Inches(0.14), top2+Inches(0.92), tw2-Inches(0.28),
            Inches(0.32), size=12, color=color)
        txt(s, desc, tx2+Inches(0.14), top2+Inches(1.28), tw2-Inches(0.28),
            Inches(3.6), size=13, color=LIGHT_GRAY)
        box(s, tx2+Inches(0.14), top2+Inches(4.98), tw2-Inches(0.28), Inches(0.34),
            fill=DARK_BG, border=color, bpt=1)
        txt(s, speed, tx2+Inches(0.18), top2+Inches(5.0), tw2-Inches(0.36),
            Inches(0.3), size=12, bold=True, color=color)
    slide_num(s, 5, T)

    # 6 — 18 Entity Categories
    s = blank(prs); bg(s); top_bar(s, A)
    txt(s, "18 Entity Categories Extracted by Layer 2 NER",
        Inches(0.6), Inches(0.18), Inches(12), Inches(0.55),
        size=26, bold=True, color=A)
    txt(s, "Original 6 expanded to 18 to support the richer Layer 3 knowledge graph.",
        Inches(0.6), Inches(0.82), Inches(12.2), Inches(0.42), size=15, color=LIGHT_GRAY)
    cats = [
        ("taxon",              "Firmicutes · Akkermansia muciniphila",         A),
        ("disease",            "IBD · T2D · colorectal cancer",               ACCENT2),
        ("method",             "16S rRNA · QIIME2 · shotgun metagenomics",    ACCENT3),
        ("body_site",          "gut · colon · oral cavity · skin",            A),
        ("treatment",          "probiotics · FMT · metformin",                ACCENT2),
        ("dataset",            "PRJNA123456 · HMP · SRA",                     ACCENT3),
        ("metabolite",         "butyrate · SCFA · bile acids · TMAO",        A),
        ("gene",               "TLR4 · NOD2 · IL-6 · FoxP3",                ACCENT2),
        ("protein",            "zonulin · calprotectin · CRP · mucin",        ACCENT3),
        ("biomarker",          "Shannon diversity · Chao1 · Bray-Curtis",     A),
        ("pathway",            "NF-κB · TLR signaling · JAK-STAT",           ACCENT2),
        ("population",         "IBD patients · healthy adults · infants",     ACCENT3),
        ("dietary_component",  "dietary fiber · inulin · polyphenols",        A),
        ("immune_cell",        "Treg · Th17 · dendritic cells · ILC3",        ACCENT2),
        ("clinical_outcome",   "remission · relapse · dysbiosis",             ACCENT3),
        ("environmental_factor","antibiotic exposure · C-section · birth mode",A),
        ("sequencing_platform","Illumina MiSeq · Oxford Nanopore · PacBio",  ACCENT2),
        ("omics_feature",      "OTU · ASV · MAG · KEGG · relative abundance", ACCENT3),
    ]
    cw2 = Inches(4.08); ch2 = Inches(0.6)
    for i, (cat, ex, color) in enumerate(cats):
        col = i % 3; row = i // 3
        cx2 = Inches(0.5) + col*(cw2+Inches(0.15))
        cy2 = Inches(1.45) + row*(ch2+Inches(0.06))
        box(s, cx2, cy2, cw2, ch2, fill=CARD_BG, border=color, bpt=1)
        txt(s, cat, cx2+Inches(0.1), cy2+Inches(0.04), Inches(1.5),
            Inches(0.28), size=12, bold=True, color=color)
        txt(s, ex,  cx2+Inches(0.1), cy2+Inches(0.32), cw2-Inches(0.2),
            Inches(0.24), size=11, color=LIGHT_GRAY)
    slide_num(s, 6, T)

    # 7 — Quality Scorer
    s = blank(prs); bg(s); top_bar(s, A)
    txt(s, "Module 8: Quality Scorer — Composite 0.0–1.0",
        Inches(0.6), Inches(0.18), Inches(12), Inches(0.55),
        size=26, bold=True, color=A)
    components = [
        ("Journal Quartile", "max 0.30",
         "Q1 → 0.30  |  Q2 → 0.20\nQ3 → 0.10  |  Q4 → 0.05\nunknown → 0.00", A),
        ("Study Design",     "max 0.30",
         "conf × 0.30\nRCT 0.95 → 0.285\nMeta-analysis → 0.22\nprotocol / letter → 0.00", ACCENT2),
        ("Sample Size",      "max 0.20",
         "≥1 000 → 0.20  |  ≥500 → 0.15\n≥100 → 0.10   |  ≥10 → 0.05\n<10 or unknown → 0.00", ACCENT3),
        ("Data Availability","max 0.15",
         "accession_linked → 0.15\nopen → 0.10\nrestricted → 0.05\nnot_stated → 0.00", LIGHT_GRAY),
        ("Sequencing Depth", "max 0.05",
         "Shotgun / WGS / metatranscriptomics → 0.05\n16S rRNA / amplicon → 0.03\nother → 0.01", A),
    ]
    cw3 = Inches(2.42)
    for i, (comp, mx, desc, color) in enumerate(components):
        cx3 = Inches(0.5) + i*(cw3+Inches(0.11))
        box(s, cx3, Inches(1.0), cw3, Inches(5.5), fill=CARD_BG, border=color, bpt=2)
        txt(s, comp, cx3+Inches(0.12), Inches(1.08), cw3-Inches(0.24),
            Inches(0.5), size=13, bold=True, color=color)
        txt(s, mx,   cx3+Inches(0.12), Inches(1.62), cw3-Inches(0.24),
            Inches(0.32), size=12, bold=True, color=WHITE)
        txt(s, desc, cx3+Inches(0.12), Inches(2.0), cw3-Inches(0.24),
            Inches(4.2), size=12, color=LIGHT_GRAY)
    txt(s, "Q1 RCT n=500 open SRA shotgun → ~0.93  |  Q2 cohort n=100 restricted 16S → ~0.60  |  Q3 review no data → ~0.30",
        Inches(0.6), Inches(6.75), Inches(12.2), Inches(0.32),
        size=12, color=ACCENT2)
    slide_num(s, 7, T)

    # 8 — EnrichedPaperRecord output
    s = blank(prs); bg(s); top_bar(s, A)
    txt(s, "Output — EnrichedPaperRecord (30+ fields)",
        Inches(0.6), Inches(0.18), Inches(12), Inches(0.55),
        size=26, bold=True, color=A)
    txt(s, "Extends PaperRecord (all Layer 1 fields) with annotations from every NLP module.",
        Inches(0.6), Inches(0.82), Inches(12.2), Inches(0.42), size=15, color=LIGHT_GRAY)
    groups2 = [
        ("Inherited (Layer 1)",  "doi · pmid · pmcid · title · abstract · mesh_terms",  A),
        ("Article Classifier",   "article_type_normalized  ·  article_type_confidence", ACCENT2),
        ("Journal Classifier",   "journal_info: IF · quartile · field · is_oa",         ACCENT3),
        ("Section Parser",       "sections: [{section_type, header, content}] × 26",   A),
        ("NER — flat",           "entities: [{text, label, ontology_id, grounded, confidence}]", ACCENT2),
        ("NER — grouped",        "taxa · diseases · methods · body_sites … (18 lists)", ACCENT3),
        ("Data Availability",    "status · accession_numbers · repositories · urls",    A),
        ("Study Design",         "type · confidence · is_rct · is_prospective",         ACCENT2),
        ("Evidence",             "sample_size · sequencing_methods · datasets",         ACCENT3),
        ("Quality & Fetch",      "quality_score 0–1  ·  fetch_source  ·  fulltext_path",A),
    ]
    gw2 = Inches(5.9); gh2 = Inches(0.78)
    for i, (grp, desc, color) in enumerate(groups2):
        col = i % 2; row = i // 2
        gx2 = Inches(0.5) + col*(gw2+Inches(0.38))
        gy2 = Inches(1.45) + row*(gh2+Inches(0.08))
        box(s, gx2, gy2, gw2, gh2, fill=CARD_BG, border=color, bpt=1)
        txt(s, grp,  gx2+Inches(0.12), gy2+Inches(0.06), Inches(2.2),
            Inches(0.3), size=12, bold=True, color=color)
        txt(s, desc, gx2+Inches(0.12), gy2+Inches(0.4), gw2-Inches(0.24),
            Inches(0.34), size=11, color=LIGHT_GRAY)
    slide_num(s, 8, T)
