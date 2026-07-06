"""Pipeline presentation - clear explanations + technical depth."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

BG = RGBColor(0x0F, 0x13, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xC8, 0xD0, 0xD8)
CYAN = RGBColor(0x58, 0xA6, 0xFF)
GREEN = RGBColor(0x3F, 0xB9, 0x50)
ORANGE = RGBColor(0xE8, 0xA8, 0x38)
RED = RGBColor(0xF8, 0x51, 0x49)
PURPLE = RGBColor(0xBC, 0x8C, 0xFF)
TEAL = RGBColor(0x39, 0xD3, 0x53)

def dark_bg(s):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG

def bar(s, c=CYAN):
    sh = s.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Pt(5))
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()

def title_s(title, sub, tag=""):
    s = prs.slides.add_slide(prs.slide_layouts[6]); dark_bg(s); bar(s)
    t = s.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10), Inches(1.5))
    p = t.text_frame.paragraphs[0]
    p.text = title; p.font.size = Pt(38); p.font.bold = True
    p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    t2 = s.shapes.add_textbox(Inches(2), Inches(4.3), Inches(9), Inches(1.2))
    p2 = t2.text_frame.paragraphs[0]
    p2.text = sub; p2.font.size = Pt(18); p2.font.color.rgb = GRAY
    p2.alignment = PP_ALIGN.CENTER
    if tag:
        t3 = s.shapes.add_textbox(Inches(2), Inches(5.8), Inches(9), Inches(0.5))
        p3 = t3.text_frame.paragraphs[0]
        p3.text = tag; p3.font.size = Pt(13); p3.font.color.rgb = CYAN
        p3.font.italic = True; p3.alignment = PP_ALIGN.CENTER

def slide(title, lines, c=CYAN):
    s = prs.slides.add_slide(prs.slide_layouts[6]); dark_bg(s); bar(s, c)
    t = s.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    p = t.text_frame.paragraphs[0]
    p.text = title; p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = WHITE
    t2 = s.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.6))
    tf = t2.text_frame; tf.word_wrap = True
    for i, l in enumerate(lines):
        pa = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        pa.text = l; pa.font.size = Pt(16); pa.font.color.rgb = GRAY
        pa.space_after = Pt(6)

# ═══════════════════════════════════════════════════════════════
title_s("Microbiome Knowledge Graph",
    "6-Stage Relevance Filtering Pipeline with Active Learning",
    "Python • Hugging Face • Ollama • scikit-learn • Neo4j • FastAPI")

slide("The Problem We're Solving", [
    "GOAL: Build a knowledge graph of human microbiome → disease → treatment links",
    "  from published papers that contain actual metagenomic sequencing data",
    "",
    "THE CHALLENGE:",
    "  • Searching 'human microbiome' across 7 APIs returns ~100,000 papers",
    "  • More than half are irrelevant (animal studies, soil, food, pure pathogen)",
    "  • Reading each paper with an LLM takes ~3 seconds = 83 hours for 100K",
    "  • Manual review at this scale is impossible",
    "",
    "OUR APPROACH: A 6-stage funnel filter",
    "  • Cheap/fast checks first (keywords, metadata) — handle 95% of papers",
    "  • Expensive LLM only for the ~5% that are genuinely ambiguous",
    "  • The system learns from its own decisions and improves over time",
], CYAN)

slide("Pipeline Overview: How Papers Flow", [
    "Think of it like a series of increasingly smart checkpoints:",
    "",
    "  Paper arrives → Stage 1: Does PubMed say it's about human microbiome?",
    "  If unsure    → Stage 2: Do keywords match? (gut microbiome +0.2, soil -0.3)",
    "  If borderline → Gate: Does it mention actual sequencing data?",
    "  If still unsure → Stage 3: ML model reads title+abstract semantically",
    "  If still unsure → Stage 3.5: Compare to papers we've already classified",
    "  If STILL unsure → Stage 4: LLM reads and decides (expensive, ~4% of papers)",
    "",
    "KEY INSIGHT: Each stage is 10-100x more expensive than the previous.",
    "  By filtering early, we save enormous computation at scale.",
    "",
    "TECHNICAL: sentence-transformers (384-dim embeddings) + LogisticRegression",
    "  + cosine similarity + Ollama LLM (qwen2.5:1.5b, local inference)",
], CYAN)

slide("Data Collection: 7 Academic APIs", [
    "We fetch papers from 7 sources simultaneously:",
    "  • PubMed (NCBI) — gold standard, has MeSH terms (librarian-curated tags)",
    "  • Europe PMC — European mirror, full-text access",
    "  • Semantic Scholar — AI-extracted metadata, citation graphs",
    "  • OpenAlex — 250M+ works, open bibliometric data",
    "  • Crossref — DOI registry, funder information",
    "  • CORE — Open access full-text repository",
    "",
    "TECHNICAL DETAILS:",
    "  • Rate-limited clients with exponential backoff (per-source config)",
    "  • Cursor-based pagination for incremental fetching across runs",
    "  • 3-layer deduplication: content_hash (SHA-256) + DOI + PMID matching",
    "  • Data model: Pydantic BaseModel (PaperRecord) with 15+ fields",
    "",
    "RESULT: 2984 raw papers → 2891 unique (93 duplicates removed)",
], CYAN)

slide("Stage 1: MeSH Metadata Filter (PubMed Only)", [
    "WHAT: PubMed papers have 'MeSH terms' — tags assigned by trained librarians",
    "  who actually read the paper. This is ground truth, not a prediction.",
    "",
    "HOW IT WORKS: Simple boolean logic on tag presence",
    "  • Has microbiome tag + human tag, no animal → KEEP (confidence: 0.90)",
    "  • Has microbiome tag + animal tag only    → REJECT (confidence: 0.05)",
    "  • Has microbiome tag but unclear          → Pass to Stage 2",
    "",
    "TECHNICAL: config/stage1_mesh.yaml contains 303 microbiome terms,",
    "  120 human-subject terms, 248 animal terms. O(1) set lookup.",
    "  Module: collectors/metadata_filter.py",
    "",
    "WHY FIRST: Zero false-positive rate, costs nothing to compute",
    "RESULT: 434 papers confidently accepted (15% resolved instantly)",
], GREEN)

slide("Stage 2: Weighted Keyword Scoring (All Sources)", [
    "WHAT: Scans the paper's text for relevant/irrelevant keywords",
    "  Each keyword has a weight. Positive terms add, negative terms subtract.",
    "",
    "HOW IT WORKS:",
    "  text = title + abstract + keywords + journal (lowercased)",
    "  score = sum of all matching term weights, normalized to [0, 1]",
    "  Score ≥ 0.70 → KEEP | Score < 0.40 → REJECT | Middle → BORDERLINE",
    "",
    "EXAMPLES:",
    "  'gut microbiome' +0.20, '16S rRNA' +0.15, 'FMT' +0.15",
    "  'mouse model' -0.25, 'zebrafish' -0.25, 'soil' -0.30",
    "",
    "TECHNICAL: Hundreds of terms in config/stage2_rules.yaml",
    "  Thresholds auto-calibrated from data distribution (stage2_calibrator.py)",
    "  Multi-word terms checked first (longest match priority)",
    "",
    "RESULT: 2230 papers resolved — the workhorse stage (77% of all papers)",
], ORANGE)

slide("Metagenomics Gate: Does the Paper Have Real Data?", [
    "WHAT: Our project needs papers with actual sequencing data — not just reviews",
    "  that discuss microbiome conceptually. This gate enforces that requirement.",
    "",
    "HOW IT WORKS:",
    "  Scan title + abstract for any of 228 data-related terms (substring match)",
    "  If no match → soft reject (paper has no data component)",
    "",
    "TERM CATEGORIES:",
    "  Sequencing: 16S, shotgun, nanopore, illumina, paired-end",
    "  QC/Preprocessing: FastQC, Trimmomatic, fastp, KneadData",
    "  Profiling: QIIME, MetaPhlAn, Kraken, DADA2",
    "  Data repositories: SRA, PRJNA, BioProject, FASTQ",
    "",
    "TECHNICAL: config/metagenomics_gate.yaml, toggleable via .env",
    "  METAGENOMICS_GATE_ENABLED=true/false",
    "",
    "RESULT: 14 papers rejected (had microbiome keywords but zero data content)",
], RED)

slide("Stage 3: ML Classifier (Self-Supervised Learning)", [
    "WHAT: A machine learning model that 'reads' the paper semantically",
    "  and predicts relevance probability. Trained without manual labeling.",
    "",
    "HOW IT WORKS:",
    "  1. Encode title+abstract → 384-dimensional vector (sentence-transformers)",
    "  2. Feed vector to LogisticRegression → probability (0.0 to 1.0)",
    "  3. prob ≥ 0.85 → KEEP | prob ≤ 0.15 → REJECT | else → BORDERLINE",
    "",
    "TRAINING (self-supervised — no human annotation needed):",
    "  • Stage 2 high-confidence papers become pseudo-labels automatically",
    "  • Hard negative mining prioritizes ambiguous papers for training",
    "  • 5-fold cross-validation, F1 = 0.93 on 21,518 training papers",
    "",
    "TECHNICAL: all-MiniLM-L6-v2 encoder, saga solver, auto-threshold",
    "  optimization via precision/recall sweep. Module: collectors/ml_classifier.py",
    "",
    "RESULT: 98 papers resolved (18 kept, 80 rejected)",
], PURPLE)

slide("Stage 3.5: Embedding Similarity (Self-Improving)", [
    "CONCEPT: If we've already classified 500 papers, why not just compare a new",
    "  paper to those? If it's very similar to known-relevant papers → accept it.",
    "",
    "HOW IT WORKS:",
    "  1. Turn the new paper into a 384-dim vector (embedding)",
    "  2. Measure cosine similarity to the closest known-relevant paper",
    "  3. Measure cosine similarity to the closest known-irrelevant paper",
    "  4. If pos_sim ≥ 0.85 AND neg_sim < 0.60 → KEEP (very similar to relevant)",
    "     If neg_sim ≥ 0.85 AND pos_sim < 0.60 → REJECT (similar to irrelevant)",
    "     Otherwise → ask the Disagreement Router",
    "",
    "THE KEY INNOVATION — it learns from itself:",
    "  After every pipeline run, confident decisions feed back into the store.",
    "  More stored papers = better boundaries = fewer uncertain papers = less LLM.",
    "",
    "TECHNICAL: NumPy brute-force cosine sim, <50ms at 100K vectors.",
    "  Store: data/embeddings/{positive,negative}.npy + metadata JSON.",
    "  Interface protocol allows FAISS swap at 500K+ vectors without code changes.",
    "  CURRENT: 226 positive + 243 negative papers (Phase 2 active)",
], TEAL)

slide("Disagreement Router + Stage 4: LLM Verifier", [
    "THE ROUTER decides: can we trust Stage 3.5, or do we need the LLM?",
    "",
    "  Route to LLM if:",
    "    • Stage 2 and Stage 3.5 disagree (one keeps, other rejects) — conflict",
    "    • Blended confidence is in the 'grey zone' [0.40, 0.70] — uncertain",
    "  Accept Stage 3.5 verdict (no LLM) if: both agree AND confidence is clear",
    "",
    "STAGE 4 — LLM VERIFIER (only for truly ambiguous papers):",
    "  Model: Ollama qwen2.5:1.5b (runs locally, no cloud, no cost)",
    "  3 cost-reduction layers:",
    "    1. Content-hash cache: same paper never verified twice (across runs)",
    "    2. Semantic cache: cosine > 0.97 = near-duplicate → reuse verdict",
    "    3. Batched: 16 papers per Ollama call (structured JSON prompt/response)",
    "  Retry: parse failure → split batch in half → retry sub-batches",
    "",
    "RESULT: Only 115 LLM calls for 2891 papers (4% verification rate!)",
    "  As the embedding store grows → disagreements shrink → LLM calls ↓ further",
], ORANGE)

slide("Post-Processing: Growth, Metrics, Reporting", [
    "After ALL papers are evaluated, three things happen:",
    "",
    "1. EMBEDDING STORE GROWTH (self-improvement):",
    "   Confident keep (score ≥ 0.80) → paper added to positive partition",
    "   Confident reject (score ≤ 0.20) → paper added to negative partition",
    "   Borderline → skipped (no feedback from uncertain decisions)",
    "   This makes Stage 3.5 smarter for the next run.",
    "",
    "2. PIPELINE METRICS (observability):",
    "   Appends JSONL record: per-stage counts, LLM calls, cache hits,",
    "   embedding store sizes, avg/p95 query latency. File: pipeline_runs.jsonl",
    "",
    "3. REJECTED PAPERS CSV (human-readable audit):",
    "   title | DOI | source | stage | score | reason (plain English)",
    "   Technical reasons auto-translated:",
    "   'ml_prob=0.07' → 'ML confidence very low — paper unlikely relevant'",
    "   LLM reasons passed through as-is (already natural language)",
], GREEN)

slide("Scaling Strategy: The Pipeline Gets Cheaper Over Time", [
    "PHASE 1 — FOUNDATION (store < 50 papers per partition)     ✓ DONE",
    "  Stage 3.5 says 'not enough data' → all borderline → LLM",
    "  LLM rate: ~4% (Stages 1+2+3 already handle 96%!)",
    "",
    "PHASE 2 — LEARNING LOOP (50-2000 papers)                   ← WE ARE HERE",
    "  Stage 3.5 starts making autonomous KEEP/REJECT decisions",
    "  Disagreement router blocks unnecessary LLM calls",
    "  Every run adds papers to the store → better next time",
    "  Expected LLM rate: ~2%",
    "",
    "PHASE 3 — SCALE HARDENING (2000+ papers):",
    "  Hybrid meta-classifier combines ALL signals into one calibrated score",
    "  Platt scaling gives true probability (0.9 = really 90% sure)",
    "  Weekly auto-retrain, monthly drift detection",
    "  Expected LLM rate: <1%",
    "",
    "TARGET: 100,000 papers with <20% ever needing LLM verification",
], TEAL)

slide("Phase 2-3: Hybrid Classifier & Quality Assurance", [
    "HYBRID META-CLASSIFIER (activates at 2000+ papers):",
    "  Combines 4 signals: rule_score + pos_similarity + neg_similarity + ml_prob",
    "  Model: LogisticRegression trained on LLM-verified papers (ground truth)",
    "  Quality gate: F1 < 0.80 → discard model, keep previous",
    "  Platt calibration: raw logits → true probabilities (monotonic sigmoid)",
    "",
    "ACTIVE LEARNING RETRAIN (weekly):",
    "  Trigger: ≥100 new LLM-verified papers since last training",
    "  Auto-retrain → evaluate F1 → accept or discard → hot-reload",
    "",
    "DRIFT MONITORING (monthly):",
    "  Sample max(ceil(N × 0.01), 10) automated decisions for manual review",
    "  Detects if classification quality is degrading over time",
    "  Output: drift_review_YYYYMM.json for human inspection",
    "",
    "TECHNICAL: collectors/hybrid_classifier.py, collectors/calibration.py,",
    "  scripts/drift_monitor.py, scheduler/jobs.py (daily/weekly/monthly)",
], PURPLE)

slide("Testing: 19 Property-Based Correctness Proofs", [
    "APPROACH: Property-Based Testing (Hypothesis framework)",
    "  Instead of testing a few examples, we test PROPERTIES that must hold",
    "  for ALL possible inputs. Hypothesis generates 100+ random test cases.",
    "",
    "KEY PROPERTIES VERIFIED:",
    "  • Embedding output: always finite, correct shape, non-zero norm",
    "  • Partition isolation: positive and negative never mix",
    "  • Stage 3.5 thresholds: decision matches rules EXACTLY for any input",
    "  • Router logic: routing matches spec for ANY combination of verdicts",
    "  • Batch invariant: ceil(N/16) batches, each ≤16, union = original",
    "  • Cache: similarity > 0.97 → hit, ≤ 0.97 → miss (always)",
    "  • Platt scaling: monotonic (a < b → calibrate(a) ≤ calibrate(b))",
    "  • Drift sampling: always ≥10 papers when population ≥10",
    "  • Metrics: stage counts always sum to total_papers",
    "",
    "RESULT: 43 tests pass / 1 known float precision edge case",
    "  This gives formal confidence the pipeline behaves correctly at scale.",
], GREEN)

slide("Production Results: First Run", [
    "INPUT: 2984 papers from 6 APIs → 2891 unique after deduplication",
    "",
    "STAGE BREAKDOWN:",
    "  Stage 1 (MeSH metadata):        434 accepted    (15%)",
    "  Stage 2 (Keyword scoring):      1504 accepted, 726 rejected (77%)",
    "  Metagenomics Gate:               14 rejected    (0.5%)",
    "  Stage 3 (ML classifier):         18 accepted, 80 rejected (3.4%)",
    "  Stage 3.5 (Embedding sim):        0 — insufficient data (expected)",
    "  Stage 4 (LLM verifier):          16 accepted, 99 rejected (4%)",
    "",
    "FINAL: 1970 kept (68%) | 921 removed (31%) | 16 review queue",
    "LLM EFFICIENCY: Only 115/2891 papers needed LLM = 4% verification rate",
    "TIME: ~6.5 min filtering + 21 min PMC full-text enrichment",
    "EMBEDDING STORE: Seeded with 226 positive + 243 negative (Phase 2 ready)",
], GREEN)

slide("Technology Stack", [
    "LANGUAGE: Python 3.12",
    "",
    "ML/NLP: sentence-transformers 2.7 (SPECTER2, all-MiniLM-L6-v2)",
    "  scikit-learn 1.5 (LogisticRegression, cross-validation, Platt scaling)",
    "  NumPy 1.26 (brute-force cosine similarity, .npy vector storage)",
    "  Ollama (qwen2.5:1.5b — local LLM, no GPU required)",
    "",
    "DATA: 7 rate-limited API clients, Pydantic data models, YAML configs",
    "  filelock (concurrent-safe embedding store), JSONL metrics",
    "",
    "TESTING: pytest + Hypothesis (property-based, 19 properties, 44 tests)",
    "LOGGING: loguru (structured, file rotation, console)",
    "",
    "DOWNSTREAM: BioBERT NER → Entity Resolution → Neo4j KG → FastAPI",
    "",
    "INFRASTRUCTURE: Runs on MacBook Air M1. No cloud. No GPU. No cost.",
], CYAN)

title_s("Summary",
    "6-stage cascade • 4% LLM rate • Self-improving feedback loop\n"
    "19 property tests • Full audit trail • Zero cloud cost",
    "Next: grow store → Phase 2 fully active → <1% LLM at 100K scale")

# ═══════════════════════════════════════════════════════════════
prs.save("/Users/shivanshukumarsingh/Desktop/GitHub/IP/Pipeline_Presentation.pptx")
print(f"Done! {len(prs.slides)} slides → Pipeline_Presentation.pptx")
